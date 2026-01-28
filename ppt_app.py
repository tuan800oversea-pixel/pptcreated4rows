import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
from PIL import Image

st.title("🖼️ 智能图片转 PPT (自动分页版)")

# --- 配置参数 ---
SLIDE_WIDTH = Inches(13.333)  # 16:9 比例
SLIDE_HEIGHT = Inches(7.5)
# 缩小顶部留白：从 1.2 英寸 缩小到 0.6 英寸
TITLE_HEIGHT = Inches(0.6)   
MARGIN = Inches(0.2)         # 左右边缘留白
SPACING = Inches(0.1)        # 图片间距
ROW_COUNT = 4                # 每页固定四行

# 初始化状态，防止下载后按钮消失
if 'ppt_data' not in st.session_state:
    st.session_state.ppt_data = None

uploaded_files = st.file_uploader("选择并上传图片 (可多选)", type=['png', 'jpg', 'jpeg'], accept_multiple_files=True)

if uploaded_files:
    # 排序图片
    files = sorted(uploaded_files, key=lambda x: x.name)
    
    if st.button("🚀 开始自动排版"):
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        
        # 计算每一行的高度
        available_height = SLIDE_HEIGHT - TITLE_HEIGHT - (2 * MARGIN) - ((ROW_COUNT - 1) * SPACING)
        row_height = available_height / ROW_COUNT
        
        def add_new_slide(p):
            return p.slides.add_slide(p.slide_layouts[6])

        # 初始化第一页
        current_slide = add_new_slide(prs)
        current_y = TITLE_HEIGHT + MARGIN
        current_x = MARGIN
        current_row = 1

        for file in files:
            img_data = Image.open(file)
            orig_w, orig_h = img_data.size
            aspect_ratio = orig_w / orig_h
            display_width = row_height * aspect_ratio
            
            # --- 换行检测 ---
            # 如果当前宽度 + 图片宽度 > 幻灯片总宽 - 右边距
            if current_x + display_width > SLIDE_WIDTH - MARGIN:
                current_x = MARGIN
                current_y += row_height + SPACING
                current_row += 1
                
                # --- 翻页检测 ---
                # 如果当前行数超过了预设的 4 行
                if current_row > ROW_COUNT:
                    current_slide = add_new_slide(prs)
                    current_y = TITLE_HEIGHT + MARGIN
                    current_row = 1
            
            # 插入图片
            current_slide.shapes.add_picture(file, current_x, current_y, height=row_height)
            
            # 移动 X 坐标
            current_x += display_width + SPACING

        # 保存结果到内存
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        st.session_state.ppt_data = ppt_buffer.getvalue()
        st.success("✅ PPT 生成成功！")

# --- 显示下载按钮 ---
if st.session_state.ppt_data:
    st.download_button(
        label="📥 点击下载 PPT 文件",
        data=st.session_state.ppt_data,
        file_name="auto_layout_presentation.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        key="download_btn" # 固定 key 确保按钮持久
    )
